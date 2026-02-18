#!/usr/bin/env python3
"""
Fetch latest investor presentation files and extract full text.

This script is optimized for mining issuers, but can work with any public company
that publishes deck files on investor relations pages.
"""

from __future__ import annotations

import argparse
import csv
import json
import logging
import re
import sys
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any
from urllib.parse import parse_qs, quote_plus, unquote, urljoin, urlparse, urlunparse

try:
    import requests
except ImportError as exc:  # pragma: no cover
    raise SystemExit(
        "Missing dependency 'requests'. Install with: pip install requests"
    ) from exc

try:
    from bs4 import BeautifulSoup
except ImportError as exc:  # pragma: no cover
    raise SystemExit(
        "Missing dependency 'beautifulsoup4'. Install with: pip install beautifulsoup4"
    ) from exc

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

try:
    from playwright.sync_api import sync_playwright
except ImportError:  # pragma: no cover
    sync_playwright = None


LOGGER = logging.getLogger("presentation_fetcher")
DOC_EXTENSIONS = {".pdf", ".ppt", ".pptx"}
US_COUNTRY_CODES = {"US", "USA", "UNITED STATES"}
DEFAULT_USER_AGENT = "mining-presentation-fetcher/1.0 (contact: replace@example.com)"
LINK_KEYWORDS = (
    "presentation",
    "investor",
    "deck",
    "corporate",
    "factsheet",
)
NAV_HINTS = (
    "events",
    "presentations",
    "investor",
    "news",
)
COMMON_IR_PATHS = (
    "/investors",
    "/investor",
    "/investor-relations",
    "/events",
    "/events-presentations",
    "/events-and-presentations",
    "/investor-relations/events-and-presentations",
    "/investors/events-and-presentations",
    "/investors/presentations",
    "/presentations",
)
EXCHANGE_ALIASES = {
    "NYSE": "NYSE",
    "NYQ": "NYSE",
    "NASDAQ": "NASDAQ",
    "NMS": "NASDAQ",
    "NASDAQGS": "NASDAQ",
    "NASDAQGM": "NASDAQ",
    "NASDAQCM": "NASDAQ",
    "AMEX": "AMEX",
    "ASE": "AMEX",
    "TSX": "TSX",
    "TOR": "TSX",
    "TSE": "TSX",
    "TSXV": "TSXV",
    "TSX-V": "TSXV",
    "TSXVENTURE": "TSXV",
    "CVE": "TSXV",
    "VAN": "TSXV",
    "CSE": "CSE",
    "CNQ": "CSE",
    "ASX": "ASX",
    "LSE": "LSE",
}
EXCHANGE_TO_COUNTRY = {
    "NYSE": "US",
    "NASDAQ": "US",
    "AMEX": "US",
    "TSX": "CA",
    "TSXV": "CA",
    "CSE": "CA",
    "ASX": "AU",
    "LSE": "GB",
}
YAHOO_SUFFIX_BY_EXCHANGE = {
    "TSX": ".TO",
    "TSXV": ".V",
    "ASX": ".AX",
    "LSE": ".L",
}


@dataclass
class Company:
    name: str
    ir_url: str | None = None
    ticker: str | None = None
    country: str | None = None
    cik: str | None = None
    exchange: str | None = None

    @property
    def id_label(self) -> str:
        if self.ticker:
            return self.ticker
        return self.name


@dataclass
class Candidate:
    company: Company
    url: str
    title: str
    source: str
    source_page: str
    published_date: datetime | None
    score: int


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Download latest investor presentation files and save full extracted text."
        )
    )
    parser.add_argument(
        "--company",
        action="append",
        default=[],
        help=(
            "Inline company spec. Use 'ticker|exchange' (example: AEM|TSX) "
            "or full form 'name|ir_url|ticker|country|cik|exchange'. "
            "Provide multiple times for multiple issuers."
        ),
    )
    parser.add_argument(
        "--companies-csv",
        type=Path,
        help="CSV file with columns: name, ir_url, ticker, country, cik, exchange",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=Path("output"),
        help="Output directory (default: ./output)",
    )
    parser.add_argument(
        "--latest-count",
        type=int,
        default=1,
        help="How many latest files to keep per company (default: 1)",
    )
    parser.add_argument(
        "--include-sec",
        action="store_true",
        help="Query SEC filing indexes as an additional US fallback.",
    )
    parser.add_argument(
        "--max-pages",
        type=int,
        default=12,
        help="Maximum IR pages to crawl per company (default: 12)",
    )
    parser.add_argument(
        "--max-filings",
        type=int,
        default=40,
        help="Maximum recent SEC filings to inspect per US company (default: 40)",
    )
    parser.add_argument(
        "--timeout",
        type=int,
        default=20,
        help="HTTP timeout in seconds (default: 20)",
    )
    parser.add_argument(
        "--user-agent",
        default=DEFAULT_USER_AGENT,
        help="User-Agent for HTTP requests, required by SEC endpoints.",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Discover and rank files without downloading them.",
    )
    parser.add_argument(
        "--disable-search-fallback",
        action="store_true",
        help="Disable DuckDuckGo-based fallback when direct crawling finds nothing.",
    )
    parser.add_argument(
        "--playwright-fallback",
        action="store_true",
        help="Use Playwright browser rendering when HTTP requests return blocked pages.",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable verbose logging.",
    )
    return parser.parse_args()


def configure_logging(verbose: bool) -> None:
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(level=level, format="%(levelname)s: %(message)s")


def slugify(text: str, max_len: int = 80) -> str:
    value = re.sub(r"[^a-zA-Z0-9]+", "-", text.strip().lower()).strip("-")
    if not value:
        value = "item"
    return value[:max_len].strip("-")


def normalize_url(url: str) -> str:
    parsed = urlparse(url)
    clean = parsed._replace(fragment="")
    return urlunparse(clean)


def to_datetime(date_text: str, fmt: str) -> datetime | None:
    try:
        return datetime.strptime(date_text, fmt)
    except ValueError:
        return None


def parse_date_from_text(text: str) -> datetime | None:
    if not text:
        return None

    iso_match = re.search(r"\b(20\d{2})[-/](\d{1,2})[-/](\d{1,2})\b", text)
    if iso_match:
        year, month, day = iso_match.groups()
        return to_datetime(f"{year}-{month}-{day}", "%Y-%m-%d")

    us_match = re.search(r"\b(\d{1,2})/(\d{1,2})/(20\d{2})\b", text)
    if us_match:
        month, day, year = us_match.groups()
        return to_datetime(f"{month}/{day}/{year}", "%m/%d/%Y")

    month_name_match = re.search(
        r"\b("
        r"jan(?:uary)?|feb(?:ruary)?|mar(?:ch)?|apr(?:il)?|may|jun(?:e)?|"
        r"jul(?:y)?|aug(?:ust)?|sep(?:t(?:ember)?)?|oct(?:ober)?|"
        r"nov(?:ember)?|dec(?:ember)?)\s+\d{1,2},\s+20\d{2}\b",
        text,
        flags=re.IGNORECASE,
    )
    if month_name_match:
        value = month_name_match.group(0)
        for fmt in ("%B %d, %Y", "%b %d, %Y"):
            parsed = to_datetime(value, fmt)
            if parsed:
                return parsed

    month_name_alt = re.search(
        r"\b\d{1,2}\s+("
        r"jan(?:uary)?|feb(?:ruary)?|mar(?:ch)?|apr(?:il)?|may|jun(?:e)?|"
        r"jul(?:y)?|aug(?:ust)?|sep(?:t(?:ember)?)?|oct(?:ober)?|"
        r"nov(?:ember)?|dec(?:ember)?)\s+20\d{2}\b",
        text,
        flags=re.IGNORECASE,
    )
    if month_name_alt:
        value = month_name_alt.group(0)
        for fmt in ("%d %B %Y", "%d %b %Y"):
            parsed = to_datetime(value, fmt)
            if parsed:
                return parsed

    year_match = re.search(r"\b(20\d{2})\b", text)
    if year_match:
        return to_datetime(f"{year_match.group(1)}-01-01", "%Y-%m-%d")
    return None


def parse_date_from_url(url: str) -> datetime | None:
    path = urlparse(url).path
    parts = [p for p in path.split("/") if p]
    combined = " ".join(parts)

    full_match = re.search(r"\b(20\d{2})(\d{2})(\d{2})\b", combined)
    if full_match:
        year, month, day = full_match.groups()
        parsed = to_datetime(f"{year}-{month}-{day}", "%Y-%m-%d")
        if parsed:
            return parsed

    ymd_match = re.search(r"\b(20\d{2})[-_](\d{1,2})[-_](\d{1,2})\b", combined)
    if ymd_match:
        year, month, day = ymd_match.groups()
        parsed = to_datetime(f"{year}-{month}-{day}", "%Y-%m-%d")
        if parsed:
            return parsed

    if len(parts) >= 3:
        for idx in range(len(parts) - 2):
            year, month, day = parts[idx : idx + 3]
            if re.fullmatch(r"20\d{2}", year) and re.fullmatch(r"\d{1,2}", month) and re.fullmatch(
                r"\d{1,2}", day
            ):
                parsed = to_datetime(f"{year}-{month}-{day}", "%Y-%m-%d")
                if parsed:
                    return parsed

    return parse_date_from_text(combined)


def same_domain(url_a: str, url_b: str) -> bool:
    return urlparse(url_a).netloc.lower() == urlparse(url_b).netloc.lower()


def is_document_link(url: str) -> bool:
    path = urlparse(url).path.lower()
    return any(path.endswith(ext) for ext in DOC_EXTENSIONS)


def text_contains_keyword(text: str) -> bool:
    lower = text.lower()
    return any(keyword in lower for keyword in LINK_KEYWORDS)


def score_link(title: str, url: str) -> int:
    lower = f"{title} {url}".lower()
    score = 0
    if "presentation" in lower:
        score += 6
    if "investor" in lower:
        score += 4
    if "deck" in lower:
        score += 3
    if "corporate" in lower:
        score += 2
    if "factsheet" in lower:
        score += 2
    if "latest" in lower:
        score += 1
    if url.lower().endswith(".pdf"):
        score += 4
    if url.lower().endswith(".pptx"):
        score += 2
    if any(term in lower for term in ("transcript", "webcast", "news release")):
        score -= 4
    return score


def infer_extension(url: str, content_type: str | None) -> str:
    suffix = Path(urlparse(url).path).suffix.lower()
    if suffix in DOC_EXTENSIONS:
        return suffix
    ct = (content_type or "").lower()
    if "pdf" in ct:
        return ".pdf"
    if "presentation" in ct or "powerpoint" in ct:
        return ".pptx"
    return ".bin"


def normalize_exchange(value: str | None) -> str | None:
    if not value:
        return None
    token = re.sub(r"\s+", "", value.strip().upper())
    return EXCHANGE_ALIASES.get(token, token if token else None)


def country_from_exchange(exchange: str | None) -> str | None:
    if not exchange:
        return None
    return EXCHANGE_TO_COUNTRY.get(normalize_exchange(exchange) or "")


def yahoo_symbol_for_ticker(ticker: str, exchange: str | None) -> str:
    symbol = ticker.strip().upper()
    if not symbol:
        return symbol
    if "." in symbol:
        return symbol
    suffix = YAHOO_SUFFIX_BY_EXCHANGE.get(normalize_exchange(exchange) or "")
    if suffix:
        return f"{symbol}{suffix}"
    return symbol


def parse_inline_parts(spec: str) -> list[str]:
    if "|" in spec:
        return [part.strip() for part in spec.split("|")]
    if "," in spec:
        return [part.strip() for part in spec.split(",")]
    return [spec.strip()]


def parse_inline_company(spec: str) -> Company:
    parts = parse_inline_parts(spec)

    # Shorthand: ticker|exchange, e.g. AEM|TSX.
    if len(parts) == 2 and parts[0] and normalize_exchange(parts[1]):
        ticker = parts[0].upper()
        exchange = normalize_exchange(parts[1])
        return Company(
            name=ticker,
            ticker=ticker,
            country=country_from_exchange(exchange),
            exchange=exchange,
        )

    if len(parts) < 6:
        parts.extend([""] * (6 - len(parts)))
    name, ir_url, ticker, country, cik, exchange = parts[:6]
    ticker = ticker.upper() if ticker else None
    exchange = normalize_exchange(exchange)
    country = country.upper() if country else None
    country = country or country_from_exchange(exchange)

    if not name and not ticker:
        raise ValueError(
            "Inline company must include at least name or ticker: "
            f"{spec!r}"
        )
    return Company(
        name=name or ticker,
        ir_url=ir_url or None,
        ticker=ticker,
        country=country,
        cik=cik or None,
        exchange=exchange,
    )


def load_companies_from_csv(path: Path) -> list[Company]:
    companies: list[Company] = []
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.DictReader(handle)
        for row in reader:
            name = (row.get("name") or "").strip()
            ticker = (row.get("ticker") or "").strip().upper()
            if not name and not ticker:
                continue
            exchange = normalize_exchange((row.get("exchange") or "").strip() or None)
            country = ((row.get("country") or "").strip() or None)
            country = country.upper() if country else None
            country = country or country_from_exchange(exchange)
            companies.append(
                Company(
                    name=name or ticker,
                    ir_url=((row.get("ir_url") or "").strip() or None),
                    ticker=(ticker or None),
                    country=country,
                    cik=((row.get("cik") or "").strip() or None),
                    exchange=exchange,
                )
            )
    return companies


def load_companies(args: argparse.Namespace) -> list[Company]:
    companies: list[Company] = []
    for spec in args.company:
        companies.append(parse_inline_company(spec))
    if args.companies_csv:
        companies.extend(load_companies_from_csv(args.companies_csv))
    if not companies:
        raise ValueError("Provide --company and/or --companies-csv input.")
    return companies


def resolve_company_market_data(
    company: Company,
    session: requests.Session,
    timeout: int,
) -> Company:
    if not company.ticker:
        return company

    query_symbol = yahoo_symbol_for_ticker(company.ticker, company.exchange)
    search_url = f"https://query2.finance.yahoo.com/v1/finance/search?q={quote_plus(query_symbol)}"

    try:
        search_resp = session.get(search_url, timeout=timeout)
        search_resp.raise_for_status()
        search_payload = search_resp.json()
    except Exception as exc:  # noqa: BLE001
        LOGGER.debug("Market-data resolve failed for %s: %s", company.ticker, exc)
        return company

    quote_candidates = search_payload.get("quotes", [])
    if not quote_candidates:
        return company

    filtered = [
        item
        for item in quote_candidates
        if str(item.get("quoteType", "")).upper() == "EQUITY"
    ] or quote_candidates
    first = filtered[0]

    resolved_name = (
        first.get("longName")
        or first.get("longname")
        or first.get("shortName")
        or first.get("shortname")
        or first.get("displayName")
        or company.name
    )
    resolved_exchange = normalize_exchange(
        company.exchange
        or first.get("exchange")
        or first.get("exchDisp")
    )
    resolved_country = company.country or country_from_exchange(resolved_exchange)

    is_placeholder_name = bool(company.ticker and company.name.upper() == company.ticker.upper())
    final_name = resolved_name if is_placeholder_name else (company.name or resolved_name)

    updated = Company(
        name=final_name,
        ir_url=company.ir_url,
        ticker=company.ticker,
        country=resolved_country,
        cik=company.cik,
        exchange=resolved_exchange,
    )

    if (
        updated.name != company.name
        or updated.country != company.country
        or updated.exchange != company.exchange
    ):
        LOGGER.debug(
            "Resolved %s -> name=%s exchange=%s country=%s",
            company.ticker,
            updated.name,
            updated.exchange,
            updated.country,
        )
    return updated


def build_ir_seed_urls(ir_url: str) -> list[str]:
    base = normalize_url(ir_url)
    parsed = urlparse(base)
    root = f"{parsed.scheme}://{parsed.netloc}"

    seeds: list[str] = [base]
    seen = {base}

    if parsed.path in ("", "/"):
        for path in COMMON_IR_PATHS:
            candidate = normalize_url(root + path)
            if candidate not in seen:
                seeds.append(candidate)
                seen.add(candidate)
    else:
        for path in COMMON_IR_PATHS:
            candidate = normalize_url(root + path)
            if candidate not in seen:
                seeds.append(candidate)
                seen.add(candidate)

    return seeds


def fetch_html_playwright(url: str, timeout: int) -> tuple[str | None, str]:
    if sync_playwright is None:
        return None, url
    try:
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch(headless=True)
            page = browser.new_page()
            page.goto(url, wait_until="domcontentloaded", timeout=timeout * 1000)
            content = page.content()
            final_url = page.url
            browser.close()
        return content, final_url
    except Exception as exc:  # noqa: BLE001
        LOGGER.debug("Playwright fetch failed for %s: %s", url, exc)
        return None, url


def fetch_html(
    session: requests.Session,
    url: str,
    timeout: int,
    use_playwright: bool = False,
) -> tuple[str | None, str]:
    try:
        response = session.get(url, timeout=timeout)
        response.raise_for_status()
    except Exception as exc:  # noqa: BLE001
        LOGGER.debug("Failed to fetch page %s: %s", url, exc)
        if use_playwright:
            return fetch_html_playwright(url, timeout=timeout)
        return None, url

    content_type = response.headers.get("content-type", "").lower()
    if "text/html" not in content_type:
        if use_playwright:
            return fetch_html_playwright(url, timeout=timeout)
        return None, str(response.url)
    return response.text, str(response.url)


def extract_links_from_html(html: str, base_url: str) -> list[tuple[str, str]]:
    soup = BeautifulSoup(html, "html.parser")
    links: list[tuple[str, str]] = []
    for anchor in soup.select("a[href]"):
        href = anchor.get("href", "").strip()
        if not href:
            continue
        abs_url = normalize_url(urljoin(base_url, href))
        text = " ".join(anchor.get_text(" ", strip=True).split())
        links.append((abs_url, text))
    return links


def discover_ir_candidates(
    company: Company,
    session: requests.Session,
    timeout: int,
    max_pages: int,
    use_playwright: bool = False,
) -> list[Candidate]:
    if not company.ir_url:
        return []

    queue: list[tuple[str, int]] = [(url, 0) for url in build_ir_seed_urls(company.ir_url)]
    visited: set[str] = set()
    candidates: list[Candidate] = []

    while queue and len(visited) < max_pages:
        page_url, depth = queue.pop(0)
        if page_url in visited:
            continue
        visited.add(page_url)

        html, final_url = fetch_html(
            session,
            page_url,
            timeout=timeout,
            use_playwright=use_playwright,
        )
        if not html:
            continue

        links = extract_links_from_html(html, final_url)
        for link_url, link_text in links:
            if not link_url.startswith("http"):
                continue

            has_keyword = text_contains_keyword(link_text) or text_contains_keyword(link_url)
            if is_document_link(link_url):
                page_has_hint = text_contains_keyword(page_url)
                if not has_keyword and not page_has_hint:
                    continue
                date_value = (
                    parse_date_from_text(link_text)
                    or parse_date_from_url(link_url)
                    or parse_date_from_url(page_url)
                )
                title = link_text or Path(urlparse(link_url).path).name
                score = score_link(title, link_url)
                if page_has_hint:
                    score += 1
                candidates.append(
                    Candidate(
                        company=company,
                        url=link_url,
                        title=title,
                        source="ir",
                        source_page=page_url,
                        published_date=date_value,
                        score=score,
                    )
                )
                continue

            if depth >= 1:
                continue
            if not same_domain(final_url, link_url):
                continue
            nav_hint = any(term in link_url.lower() for term in NAV_HINTS) or any(
                term in link_text.lower() for term in NAV_HINTS
            )
            if has_keyword or nav_hint:
                queue.append((link_url, depth + 1))

    return candidates


def decode_duckduckgo_target(raw_url: str) -> str:
    parsed = urlparse(raw_url)
    if "duckduckgo.com" in parsed.netloc and parsed.path.startswith("/l/"):
        query = parse_qs(parsed.query)
        target = query.get("uddg")
        if target and target[0]:
            return normalize_url(unquote(target[0]))
    return normalize_url(raw_url)


def discover_search_candidates(
    company: Company,
    session: requests.Session,
    timeout: int,
    max_results: int = 15,
) -> list[Candidate]:
    query_parts = [company.name, "investor presentation pdf"]
    if company.ticker:
        query_parts.append(company.ticker)
    if company.exchange:
        query_parts.append(company.exchange)
    query = " ".join(part for part in query_parts if part)
    search_url = f"https://duckduckgo.com/html/?q={quote_plus(query)}"

    html, final_url = fetch_html(session, search_url, timeout=timeout)
    if not html:
        return []

    links = extract_links_from_html(html, final_url)
    candidates: list[Candidate] = []
    for link_url, link_text in links:
        if len(candidates) >= max_results:
            break
        target_url = decode_duckduckgo_target(link_url)
        if not target_url.startswith("http"):
            continue
        if "duckduckgo.com" in urlparse(target_url).netloc:
            continue

        has_keyword = text_contains_keyword(link_text) or text_contains_keyword(target_url)
        if not has_keyword and not is_document_link(target_url):
            continue

        score = score_link(link_text, target_url)
        if "duckduckgo.com" in target_url:
            continue
        if "presentation" in target_url.lower():
            score += 1

        date_value = parse_date_from_text(link_text) or parse_date_from_url(target_url)
        title = link_text or Path(urlparse(target_url).path).name or "presentation"
        candidates.append(
            Candidate(
                company=company,
                url=target_url,
                title=title,
                source="search",
                source_page=search_url,
                published_date=date_value,
                score=score,
            )
        )
    return candidates


def discover_bing_candidates(
    company: Company,
    session: requests.Session,
    timeout: int,
    max_results: int = 15,
) -> list[Candidate]:
    query_parts = [company.name, "investor presentation pdf"]
    if company.ticker:
        query_parts.append(company.ticker)
    if company.exchange:
        query_parts.append(company.exchange)
    query = " ".join(part for part in query_parts if part)
    search_url = f"https://www.bing.com/search?q={quote_plus(query)}"

    html, final_url = fetch_html(session, search_url, timeout=timeout)
    if not html:
        return []

    soup = BeautifulSoup(html, "html.parser")
    nodes = soup.select("li.b_algo h2 a[href]") or soup.select("a[href]")
    candidates: list[Candidate] = []

    for node in nodes:
        if len(candidates) >= max_results:
            break
        href = str(node.get("href", "")).strip()
        if not href:
            continue
        target_url = normalize_url(urljoin(final_url, href))
        if not target_url.startswith("http"):
            continue
        netloc = urlparse(target_url).netloc.lower()
        if "bing.com" in netloc:
            continue

        title = " ".join(node.get_text(" ", strip=True).split())
        has_keyword = text_contains_keyword(title) or text_contains_keyword(target_url)
        if not has_keyword and not is_document_link(target_url):
            continue

        score = score_link(title, target_url)
        date_value = parse_date_from_text(title) or parse_date_from_url(target_url)
        candidates.append(
            Candidate(
                company=company,
                url=target_url,
                title=title or Path(urlparse(target_url).path).name or "presentation",
                source="search",
                source_page=search_url,
                published_date=date_value,
                score=score,
            )
        )
    return candidates


def discover_web_search_candidates(
    company: Company,
    session: requests.Session,
    timeout: int,
    max_results: int = 15,
) -> list[Candidate]:
    candidates = discover_search_candidates(
        company=company,
        session=session,
        timeout=timeout,
        max_results=max_results,
    )
    if candidates:
        return candidates
    return discover_bing_candidates(
        company=company,
        session=session,
        timeout=timeout,
        max_results=max_results,
    )


def normalize_cik(raw_cik: str) -> str:
    digits = re.sub(r"\D", "", raw_cik)
    return digits.zfill(10)


def sec_ticker_to_cik(session: requests.Session, timeout: int) -> dict[str, str]:
    url = "https://www.sec.gov/files/company_tickers.json"
    response = session.get(url, timeout=timeout)
    response.raise_for_status()
    payload = response.json()
    mapping: dict[str, str] = {}
    for item in payload.values():
        ticker = str(item.get("ticker", "")).upper().strip()
        cik_str = str(item.get("cik_str", "")).strip()
        if ticker and cik_str:
            mapping[ticker] = normalize_cik(cik_str)
    return mapping


def resolve_cik(
    company: Company,
    ticker_map: dict[str, str] | None,
) -> str | None:
    if company.cik:
        return normalize_cik(company.cik)
    if company.ticker and ticker_map:
        return ticker_map.get(company.ticker.upper())
    return None


def discover_sec_candidates(
    company: Company,
    session: requests.Session,
    timeout: int,
    max_filings: int,
    ticker_map: dict[str, str] | None,
) -> list[Candidate]:
    cik = resolve_cik(company, ticker_map)
    if not cik:
        return []

    submissions_url = f"https://data.sec.gov/submissions/CIK{cik}.json"
    try:
        response = session.get(submissions_url, timeout=timeout)
        response.raise_for_status()
    except Exception as exc:  # noqa: BLE001
        LOGGER.debug("Failed SEC submissions for %s: %s", company.id_label, exc)
        return []

    payload = response.json()
    recent = payload.get("filings", {}).get("recent", {})
    forms = recent.get("form", [])
    filing_dates = recent.get("filingDate", [])
    accession_numbers = recent.get("accessionNumber", [])

    candidates: list[Candidate] = []
    scanned = 0
    for idx, form in enumerate(forms):
        if scanned >= max_filings:
            break
        scanned += 1
        if form not in {"8-K", "6-K"}:
            continue

        accession = accession_numbers[idx]
        filing_date = parse_date_from_text(filing_dates[idx] if idx < len(filing_dates) else "")
        accession_nodash = accession.replace("-", "")
        archive_base = f"https://www.sec.gov/Archives/edgar/data/{int(cik)}/{accession_nodash}/"
        index_url = archive_base + "index.json"

        try:
            index_response = session.get(index_url, timeout=timeout)
            index_response.raise_for_status()
            index_payload = index_response.json()
        except Exception as exc:  # noqa: BLE001
            LOGGER.debug("Failed SEC index %s: %s", index_url, exc)
            continue

        items = index_payload.get("directory", {}).get("item", [])
        for item in items:
            name = str(item.get("name", ""))
            lower_name = name.lower()
            if not any(lower_name.endswith(ext) for ext in DOC_EXTENSIONS):
                continue
            if not text_contains_keyword(lower_name):
                continue
            doc_url = archive_base + name
            score = score_link(name, doc_url) + 2
            candidates.append(
                Candidate(
                    company=company,
                    url=doc_url,
                    title=f"{form} {name}",
                    source="sec",
                    source_page=index_url,
                    published_date=filing_date,
                    score=score,
                )
            )
    return candidates


def dedupe_candidates(candidates: list[Candidate]) -> list[Candidate]:
    by_url: dict[str, Candidate] = {}
    for candidate in candidates:
        existing = by_url.get(candidate.url)
        if not existing:
            by_url[candidate.url] = candidate
            continue
        candidate_date = candidate.published_date or datetime(1970, 1, 1)
        existing_date = existing.published_date or datetime(1970, 1, 1)
        if candidate_date > existing_date or (
            candidate_date == existing_date and candidate.score > existing.score
        ):
            by_url[candidate.url] = candidate
    return list(by_url.values())


def pick_latest(candidates: list[Candidate], count: int) -> list[Candidate]:
    deduped = dedupe_candidates(candidates)
    sorted_candidates = sorted(
        deduped,
        key=lambda item: (
            item.published_date or datetime(1970, 1, 1),
            item.score,
        ),
        reverse=True,
    )
    return sorted_candidates[:count]


def extract_pdf_text(path: Path) -> tuple[str, str]:
    if PdfReader is None:
        return "", "pdf-skip-pypdf-missing"
    try:
        reader = PdfReader(str(path))
        chunks: list[str] = []
        for idx, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            chunks.append(f"--- page {idx} ---\n{page_text}")
        return "\n\n".join(chunks).strip(), "pdf-pypdf"
    except Exception as exc:  # noqa: BLE001
        return "", f"pdf-error:{exc}"


def extract_pptx_text(path: Path) -> tuple[str, str]:
    if Presentation is None:
        return "", "pptx-skip-python-pptx-missing"
    try:
        presentation = Presentation(str(path))
        chunks: list[str] = []
        for idx, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        slide_lines.append(text)
            chunks.append(f"--- slide {idx} ---\n" + "\n".join(slide_lines))
        return "\n\n".join(chunks).strip(), "pptx-python-pptx"
    except Exception as exc:  # noqa: BLE001
        return "", f"pptx-error:{exc}"


def extract_text_from_file(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path)
    if suffix == ".pptx":
        return extract_pptx_text(path)
    return "", f"unsupported:{suffix or 'none'}"


def download_bytes(
    session: requests.Session,
    url: str,
    timeout: int,
) -> tuple[bytes, str | None]:
    response = session.get(url, timeout=timeout)
    response.raise_for_status()
    return response.content, response.headers.get("content-type")


def save_candidate(
    candidate: Candidate,
    session: requests.Session,
    output_dir: Path,
    timeout: int,
    dry_run: bool,
) -> dict[str, Any]:
    company_dir = output_dir / slugify(candidate.company.id_label, max_len=50)
    company_dir.mkdir(parents=True, exist_ok=True)

    date_label = (
        candidate.published_date.strftime("%Y-%m-%d")
        if candidate.published_date
        else "undated"
    )
    base_name = f"{date_label}__{slugify(candidate.title, max_len=80)}"
    metadata: dict[str, Any] = {
        "company_name": candidate.company.name,
        "ticker": candidate.company.ticker,
        "exchange": candidate.company.exchange,
        "country": candidate.company.country,
        "source": candidate.source,
        "source_page": candidate.source_page,
        "document_url": candidate.url,
        "title": candidate.title,
        "published_date": candidate.published_date.strftime("%Y-%m-%d")
        if candidate.published_date
        else None,
        "score": candidate.score,
        "downloaded_at": datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ"),
        "local_file": None,
        "text_file": None,
        "text_extraction_method": None,
        "text_char_count": 0,
    }

    if dry_run:
        return metadata

    content, content_type = download_bytes(session, candidate.url, timeout=timeout)
    extension = infer_extension(candidate.url, content_type)
    file_path = company_dir / f"{base_name}{extension}"
    file_path.write_bytes(content)
    metadata["local_file"] = str(file_path)

    text, method = extract_text_from_file(file_path)
    metadata["text_extraction_method"] = method
    if text:
        text_path = company_dir / f"{base_name}.txt"
        text_path.write_text(text, encoding="utf-8")
        metadata["text_file"] = str(text_path)
        metadata["text_char_count"] = len(text)

    metadata_path = company_dir / f"{base_name}.metadata.json"
    metadata_path.write_text(json.dumps(metadata, indent=2), encoding="utf-8")
    return metadata


def write_manifest(records: list[dict[str, Any]], output_dir: Path) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    json_path = output_dir / "manifest.json"
    csv_path = output_dir / "manifest.csv"
    json_path.write_text(json.dumps(records, indent=2), encoding="utf-8")

    if not records:
        csv_path.write_text("", encoding="utf-8")
        return

    fieldnames: list[str] = []
    for record in records:
        for key in record.keys():
            if key not in fieldnames:
                fieldnames.append(key)

    with csv_path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames)
        writer.writeheader()
        for record in records:
            writer.writerow(record)


def should_query_sec(company: Company) -> bool:
    if company.cik:
        return True
    if company.exchange:
        exchange_country = country_from_exchange(company.exchange)
        if exchange_country and exchange_country not in US_COUNTRY_CODES:
            return False
    if company.country and company.country.upper() not in US_COUNTRY_CODES:
        return False
    return bool(company.ticker)


def main() -> int:
    args = parse_args()
    configure_logging(args.verbose)

    if args.latest_count < 1:
        raise SystemExit("--latest-count must be >= 1")
    if args.playwright_fallback and sync_playwright is None:
        LOGGER.warning(
            "Playwright fallback requested but playwright is not installed. "
            "Install with: pip install playwright && playwright install chromium"
        )

    try:
        companies = load_companies(args)
    except Exception as exc:  # noqa: BLE001
        raise SystemExit(f"Invalid company input: {exc}") from exc

    session = requests.Session()
    session.headers.update(
        {
            "User-Agent": args.user_agent,
            "Accept": "text/html,application/json,application/pdf,*/*",
            "Accept-Language": "en-US,en;q=0.9",
        }
    )

    companies = [
        resolve_company_market_data(
            company=company,
            session=session,
            timeout=args.timeout,
        )
        for company in companies
    ]

    ticker_map: dict[str, str] | None = None
    if args.include_sec:
        try:
            ticker_map = sec_ticker_to_cik(session, timeout=args.timeout)
        except Exception as exc:  # noqa: BLE001
            LOGGER.warning("Unable to load SEC ticker map: %s", exc)
            ticker_map = None

    all_records: list[dict[str, Any]] = []
    for company in companies:
        LOGGER.info("Processing %s", company.id_label)

        candidates: list[Candidate] = []
        candidates.extend(
            discover_ir_candidates(
                company=company,
                    session=session,
                    timeout=args.timeout,
                    max_pages=args.max_pages,
                    use_playwright=args.playwright_fallback,
                )
            )

        if args.include_sec and should_query_sec(company):
            candidates.extend(
                discover_sec_candidates(
                    company=company,
                    session=session,
                    timeout=args.timeout,
                    max_filings=args.max_filings,
                    ticker_map=ticker_map,
                )
            )

        if not candidates and not args.disable_search_fallback:
            candidates.extend(
                discover_web_search_candidates(
                    company=company,
                    session=session,
                    timeout=args.timeout,
                )
            )

        selected = pick_latest(candidates, count=args.latest_count)
        if not selected:
            LOGGER.warning("No presentation candidates found for %s", company.id_label)
            continue

        for candidate in selected:
            LOGGER.info(
                "Selected %s (%s)",
                candidate.url,
                candidate.published_date.strftime("%Y-%m-%d")
                if candidate.published_date
                else "undated",
            )
            try:
                record = save_candidate(
                    candidate=candidate,
                    session=session,
                    output_dir=args.output_dir,
                    timeout=args.timeout,
                    dry_run=args.dry_run,
                )
                all_records.append(record)
            except Exception as exc:  # noqa: BLE001
                LOGGER.error("Failed to save %s: %s", candidate.url, exc)

    write_manifest(all_records, args.output_dir)
    LOGGER.info("Completed. Records: %s. Manifest: %s", len(all_records), args.output_dir)
    return 0


if __name__ == "__main__":
    sys.exit(main())
